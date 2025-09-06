"""
Sistema de Generación de Documentos Automática
Genera documentos en Google Workspace y Microsoft Office a partir de datos de scraping
"""

import asyncio
import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

# Importaciones para generación de documentos
try:
    # Microsoft Office
    import openpyxl
    from docx import Document
    from docx.enum.style import WD_STYLE_TYPE
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt
    from openpyxl.styles import Font, PatternFill
    from pptx import Presentation

    OFFICE_AVAILABLE = True
except ImportError:
    OFFICE_AVAILABLE = False

try:
    # Google Workspace API
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    GOOGLE_API_AVAILABLE = True
except ImportError:
    GOOGLE_API_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentGenerator:
    """Generador de documentos automático"""

    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or {}
        self.output_dir = Path(self.config.get("output_dir", "exports/documents"))
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Configuración de Google API
        self.google_credentials = None
        self.google_services = {}

        # Templates y estilos
        self.templates = {
            "research_report": "plantilla_informe_investigacion",
            "data_analysis": "plantilla_analisis_datos",
            "summary": "plantilla_resumen",
            "presentation": "plantilla_presentacion",
        }

        logger.info("📄 Generador de documentos inicializado")

    def setup_google_credentials(
        self, credentials_path: str = None, service_account_path: str = None
    ):
        """Configura las credenciales de Google API"""
        try:
            if service_account_path and os.path.exists(service_account_path):
                # Usar cuenta de servicio
                self.google_credentials = (
                    service_account.Credentials.from_service_account_file(
                        service_account_path,
                        scopes=[
                            "https://www.googleapis.com/auth/documents",
                            "https://www.googleapis.com/auth/spreadsheets",
                            "https://www.googleapis.com/auth/presentations",
                            "https://www.googleapis.com/auth/drive.file",
                        ],
                    )
                )

                # Construir servicios
                self.google_services = {
                    "docs": build("docs", "v1", credentials=self.google_credentials),
                    "sheets": build(
                        "sheets", "v4", credentials=self.google_credentials
                    ),
                    "slides": build(
                        "slides", "v1", credentials=self.google_credentials
                    ),
                    "drive": build("drive", "v3", credentials=self.google_credentials),
                }

                logger.info("✅ Credenciales de Google configuradas exitosamente")
                return True

            else:
                logger.warning("⚠️ No se encontraron credenciales de Google válidas")
                return False

        except Exception as e:
            logger.error(f"Error configurando credenciales de Google: {e}")
            return False

    async def generate_comprehensive_document(
        self,
        search_data: Dict[str, Any],
        document_type: str = "research_report",
        format_types: List[str] = None,
    ) -> Dict[str, Any]:
        """Genera documento comprensivo basado en datos de búsqueda"""

        if format_types is None:
            format_types = ["docx", "md"]

        generation_session = {
            "session_id": f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "document_type": document_type,
            "format_types": format_types,
            "timestamp": datetime.now().isoformat(),
            "status": "starting",
            "generated_files": [],
        }

        try:
            topic = search_data.get("topic", "Documento sin título")
            synthesis = search_data.get("synthesis", {})
            organized_content = search_data.get("organized_content", {})

            logger.info(f"📄 Generando documento '{document_type}' para: {topic}")

            # Preparar contenido estructurado
            structured_content = self._structure_content_for_document(
                search_data, document_type
            )

            generation_session["structured_content"] = structured_content
            generation_session["status"] = "generating"

            # Generar en cada formato solicitado
            for format_type in format_types:
                try:
                    if format_type == "docx" and OFFICE_AVAILABLE:
                        file_path = await self._generate_word_document(
                            structured_content, topic, document_type
                        )
                        generation_session["generated_files"].append(
                            {"format": "docx", "path": file_path, "status": "success"}
                        )

                    elif format_type == "xlsx" and OFFICE_AVAILABLE:
                        file_path = await self._generate_excel_document(
                            structured_content, topic, document_type
                        )
                        generation_session["generated_files"].append(
                            {"format": "xlsx", "path": file_path, "status": "success"}
                        )

                    elif format_type == "pptx" and OFFICE_AVAILABLE:
                        file_path = await self._generate_powerpoint_document(
                            structured_content, topic, document_type
                        )
                        generation_session["generated_files"].append(
                            {"format": "pptx", "path": file_path, "status": "success"}
                        )

                    elif (
                        format_type == "gdoc"
                        and GOOGLE_API_AVAILABLE
                        and self.google_credentials
                    ):
                        doc_id = await self._generate_google_doc(
                            structured_content, topic, document_type
                        )
                        generation_session["generated_files"].append(
                            {
                                "format": "gdoc",
                                "id": doc_id,
                                "url": f"https://docs.google.com/document/d/{doc_id}",
                                "status": "success",
                            }
                        )

                    elif format_type == "md":
                        file_path = await self._generate_markdown_document(
                            structured_content, topic, document_type
                        )
                        generation_session["generated_files"].append(
                            {"format": "md", "path": file_path, "status": "success"}
                        )

                    else:
                        generation_session["generated_files"].append(
                            {
                                "format": format_type,
                                "status": "not_available",
                                "reason": "Dependencies not installed or credentials not configured",
                            }
                        )

                except Exception as e:
                    logger.error(f"Error generando formato {format_type}: {e}")
                    generation_session["generated_files"].append(
                        {"format": format_type, "status": "error", "error": str(e)}
                    )

            generation_session["status"] = "completed"
            generation_session["completion_time"] = datetime.now().isoformat()

            successful_files = [
                f
                for f in generation_session["generated_files"]
                if f["status"] == "success"
            ]

            logger.info(
                f"📄 Documento generado: {len(successful_files)} archivos creados"
            )

            return generation_session

        except Exception as e:
            logger.error(f"Error en generación de documento: {e}")
            generation_session["status"] = "error"
            generation_session["error"] = str(e)
            return generation_session

    def _structure_content_for_document(
        self, search_data: Dict[str, Any], document_type: str
    ) -> Dict[str, Any]:
        """Estructura el contenido según el tipo de documento"""

        topic = search_data.get("topic", "Documento")
        synthesis = search_data.get("synthesis", {})
        organized_content = search_data.get("organized_content", {})

        if document_type == "research_report":
            return {
                "title": f"Informe de Investigación: {topic}",
                "sections": [
                    {
                        "title": "Resumen Ejecutivo",
                        "content": synthesis.get("summary", "Sin resumen disponible"),
                        "type": "summary",
                    },
                    {
                        "title": "Metodología",
                        "content": f"Se realizó una búsqueda inteligente consultando {synthesis.get('sources_count', 0)} fuentes "
                        f"con una calidad de contenido de {synthesis.get('content_quality', 0):.1%}.",
                        "type": "methodology",
                    },
                    {
                        "title": "Hallazgos Principales",
                        "content": synthesis.get("key_points", []),
                        "type": "findings",
                    },
                    {
                        "title": "Contenido Detallado",
                        "content": organized_content,
                        "type": "detailed_content",
                    },
                    {
                        "title": "Conclusiones y Recomendaciones",
                        "content": synthesis.get("recommended_actions", []),
                        "type": "conclusions",
                    },
                ],
                "metadata": {
                    "author": "Web Scraper PRO",
                    "creation_date": datetime.now().strftime("%Y-%m-%d"),
                    "topic": topic,
                    "sources": synthesis.get("sources_count", 0),
                },
            }

        elif document_type == "data_analysis":
            return {
                "title": f"Análisis de Datos: {topic}",
                "sections": [
                    {
                        "title": "Objetivos del Análisis",
                        "content": f"Analizar la información disponible sobre {topic} para extraer insights significativos.",
                        "type": "objectives",
                    },
                    {
                        "title": "Datos Recopilados",
                        "content": organized_content,
                        "type": "data",
                    },
                    {
                        "title": "Análisis de Calidad",
                        "content": f"Calidad de los datos: {synthesis.get('content_quality', 0):.1%}",
                        "type": "quality_analysis",
                    },
                    {
                        "title": "Insights Clave",
                        "content": synthesis.get("key_points", []),
                        "type": "insights",
                    },
                ],
                "metadata": {
                    "analysis_type": "content_analysis",
                    "data_sources": synthesis.get("sources_count", 0),
                    "topic": topic,
                },
            }

        else:  # summary por defecto
            return {
                "title": f"Resumen: {topic}",
                "sections": [
                    {
                        "title": "Información Principal",
                        "content": synthesis.get("summary", ""),
                        "type": "main_info",
                    },
                    {
                        "title": "Puntos Destacados",
                        "content": synthesis.get("key_points", []),
                        "type": "highlights",
                    },
                ],
                "metadata": {"document_type": "summary", "topic": topic},
            }

    async def _generate_word_document(
        self, structured_content: Dict[str, Any], topic: str, document_type: str
    ) -> str:
        """Genera documento de Word"""

        doc = Document()

        # Configurar estilos
        self._setup_word_styles(doc)

        # Título principal
        title = doc.add_heading(structured_content["title"], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadatos
        metadata = structured_content.get("metadata", {})
        info_paragraph = doc.add_paragraph()
        info_paragraph.add_run(f"Autor: {metadata.get('author', 'N/A')}\n").bold = True
        info_paragraph.add_run(f"Fecha: {metadata.get('creation_date', 'N/A')}\n")
        info_paragraph.add_run(f"Fuentes consultadas: {metadata.get('sources', 'N/A')}")

        doc.add_page_break()

        # Contenido por secciones
        for section in structured_content.get("sections", []):
            # Título de sección
            doc.add_heading(section["title"], 1)

            # Contenido de sección
            if isinstance(section["content"], str):
                doc.add_paragraph(section["content"])

            elif isinstance(section["content"], list):
                for item in section["content"]:
                    if isinstance(item, str):
                        p = doc.add_paragraph(item, style="List Bullet")
                    else:
                        p = doc.add_paragraph(str(item), style="List Bullet")

            elif isinstance(section["content"], dict):
                # Contenido organizado (como organized_content)
                for content_type, results in section["content"].items():
                    if results:
                        doc.add_heading(content_type.title(), 2)

                        for result in results:
                            if isinstance(result, dict):
                                title = result.get("title", "Sin título")
                                content = result.get("content", "Sin contenido")
                                url = result.get("url", "")

                                p = doc.add_paragraph()
                                p.add_run(title).bold = True
                                p.add_paragraph(
                                    content[:500]
                                    + ("..." if len(content) > 500 else "")
                                )
                                if url:
                                    p.add_paragraph(f"Fuente: {url}")

                                doc.add_paragraph()  # Espacio

        # Guardar documento
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{document_type}_{topic.replace(' ', '_')}_{timestamp}.docx"
        filepath = self.output_dir / filename

        doc.save(str(filepath))

        logger.info(f"📄 Documento Word generado: {filepath}")
        return str(filepath)

    def _setup_word_styles(self, doc):
        """Configura estilos para el documento Word"""
        try:
            # Estilo para párrafos normales
            styles = doc.styles

            # Estilo personalizado para contenido
            if "Content" not in [style.name for style in styles]:
                content_style = styles.add_style("Content", WD_STYLE_TYPE.PARAGRAPH)
                content_style.font.name = "Arial"
                content_style.font.size = Pt(11)
        except Exception as e:
            logger.warning(f"No se pudieron configurar estilos personalizados: {e}")

    async def _generate_excel_document(
        self, structured_content: Dict[str, Any], topic: str, document_type: str
    ) -> str:
        """Genera documento de Excel"""

        wb = openpyxl.Workbook()

        # Hoja principal con resumen
        ws = wb.active
        ws.title = "Resumen"

        # Título
        ws["A1"] = structured_content["title"]
        ws["A1"].font = Font(size=16, bold=True)
        ws.merge_cells("A1:D1")

        # Metadatos
        metadata = structured_content.get("metadata", {})
        row = 3
        for key, value in metadata.items():
            ws[f"A{row}"] = key.replace("_", " ").title() + ":"
            ws[f"B{row}"] = value
            ws[f"A{row}"].font = Font(bold=True)
            row += 1

        # Secciones
        row += 2
        for section in structured_content.get("sections", []):
            ws[f"A{row}"] = section["title"]
            ws[f"A{row}"].font = Font(size=14, bold=True)
            row += 1

            if isinstance(section["content"], str):
                ws[f"A{row}"] = section["content"]
                row += 2

            elif isinstance(section["content"], list):
                for item in section["content"]:
                    ws[f"B{row}"] = f"• {item}"
                    row += 1
                row += 1

        # Ajustar ancho de columnas
        ws.column_dimensions["A"].width = 30
        ws.column_dimensions["B"].width = 60

        # Crear hoja de datos si hay contenido organizado
        for section in structured_content.get("sections", []):
            if isinstance(section["content"], dict) and section["content"]:
                self._create_excel_data_sheet(wb, section["content"], section["title"])

        # Guardar archivo
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{document_type}_{topic.replace(' ', '_')}_{timestamp}.xlsx"
        filepath = self.output_dir / filename

        wb.save(str(filepath))

        logger.info(f"📊 Documento Excel generado: {filepath}")
        return str(filepath)

    def _create_excel_data_sheet(
        self, workbook, organized_content: Dict[str, Any], sheet_name: str
    ):
        """Crea hoja de datos en Excel"""
        ws = workbook.create_sheet(
            title=sheet_name[:31]
        )  # Límite de caracteres en Excel

        # Headers
        headers = [
            "Tipo",
            "Título",
            "Fuente",
            "URL",
            "Confianza",
            "Contenido (Preview)",
        ]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(
                start_color="CCCCCC", end_color="CCCCCC", fill_type="solid"
            )

        # Datos
        row = 2
        for content_type, results in organized_content.items():
            for result in results:
                if isinstance(result, dict):
                    ws.cell(row=row, column=1, value=content_type)
                    ws.cell(row=row, column=2, value=result.get("title", ""))
                    ws.cell(row=row, column=3, value=result.get("source", ""))
                    ws.cell(row=row, column=4, value=result.get("url", ""))
                    ws.cell(
                        row=row, column=5, value=f"{result.get('confidence', 0):.1%}"
                    )

                    content_preview = result.get("content", "")[:100]
                    ws.cell(row=row, column=6, value=content_preview)

                    row += 1

        # Ajustar anchos
        for col in range(1, 7):
            ws.column_dimensions[chr(64 + col)].width = 20

    async def _generate_powerpoint_document(
        self, structured_content: Dict[str, Any], topic: str, document_type: str
    ) -> str:
        """Genera presentación de PowerPoint"""

        prs = Presentation()

        # Slide de título
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = structured_content["title"]

        metadata = structured_content.get("metadata", {})
        subtitle.text = (
            f"Generado por Web Scraper PRO\n{metadata.get('creation_date', '')}"
        )

        # Slides de contenido
        for section in structured_content.get("sections", []):
            # Layout de contenido
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Título de slide
            title = slide.shapes.title
            title.text = section["title"]

            # Contenido
            content = slide.placeholders[1]
            text_frame = content.text_frame

            if isinstance(section["content"], str):
                text_frame.text = section["content"][:500]  # Limitar texto

            elif isinstance(section["content"], list):
                text_frame.text = ""
                for item in section["content"]:
                    p = text_frame.add_paragraph()
                    p.text = str(item)[:200]  # Limitar cada punto
                    p.level = 1

        # Guardar presentación
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{document_type}_{topic.replace(' ', '_')}_{timestamp}.pptx"
        filepath = self.output_dir / filename

        prs.save(str(filepath))

        logger.info(f"📊 Presentación PowerPoint generada: {filepath}")
        return str(filepath)

    async def _generate_google_doc(
        self, structured_content: Dict[str, Any], topic: str, document_type: str
    ) -> str:
        """Genera documento de Google Docs"""

        if not self.google_services.get("docs"):
            raise Exception("Servicio de Google Docs no disponible")

        try:
            # Crear documento
            title = structured_content["title"]
            body = {"title": title}

            document = (
                self.google_services["docs"].documents().create(body=body).execute()
            )
            document_id = document.get("documentId")

            # Preparar contenido
            requests = []

            # Contenido de secciones
            content_text = ""
            for section in structured_content.get("sections", []):
                content_text += f"\n\n{section['title']}\n"
                content_text += "=" * len(section["title"]) + "\n\n"

                if isinstance(section["content"], str):
                    content_text += section["content"]
                elif isinstance(section["content"], list):
                    for item in section["content"]:
                        content_text += f"• {item}\n"

                content_text += "\n"

            # Insertar texto
            requests.append(
                {"insertText": {"location": {"index": 1}, "text": content_text}}
            )

            # Ejecutar requests
            if requests:
                self.google_services["docs"].documents().batchUpdate(
                    documentId=document_id, body={"requests": requests}
                ).execute()

            logger.info(f"📄 Google Doc generado: {document_id}")
            return document_id

        except Exception as e:
            logger.error(f"Error generando Google Doc: {e}")
            raise

    async def _generate_markdown_document(
        self, structured_content: Dict[str, Any], topic: str, document_type: str
    ) -> str:
        """Genera documento Markdown"""

        md_content = f"# {structured_content['title']}\n\n"

        # Metadatos
        metadata = structured_content.get("metadata", {})
        if metadata:
            md_content += "## Información del Documento\n\n"
            for key, value in metadata.items():
                md_content += f"**{key.replace('_', ' ').title()}:** {value}\n\n"

        # Secciones
        for section in structured_content.get("sections", []):
            md_content += f"## {section['title']}\n\n"

            if isinstance(section["content"], str):
                md_content += f"{section['content']}\n\n"

            elif isinstance(section["content"], list):
                for item in section["content"]:
                    md_content += f"- {item}\n"
                md_content += "\n"

            elif isinstance(section["content"], dict):
                for content_type, results in section["content"].items():
                    if results:
                        md_content += f"### {content_type.title()}\n\n"

                        for result in results:
                            if isinstance(result, dict):
                                title = result.get("title", "Sin título")
                                content = result.get("content", "")[:300]
                                url = result.get("url", "")

                                md_content += f"**{title}**\n\n"
                                md_content += f"{content}...\n\n"
                                if url:
                                    md_content += f"[Ver fuente]({url})\n\n"
                                md_content += "---\n\n"

        # Guardar archivo
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{document_type}_{topic.replace(' ', '_')}_{timestamp}.md"
        filepath = self.output_dir / filename

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(md_content)

        logger.info(f"📄 Documento Markdown generado: {filepath}")
        return str(filepath)

    def list_generated_documents(self) -> List[Dict[str, Any]]:
        """Lista documentos generados"""

        documents = []

        for file_path in self.output_dir.glob("*"):
            if file_path.is_file():
                documents.append(
                    {
                        "name": file_path.name,
                        "path": str(file_path),
                        "size": file_path.stat().st_size,
                        "created": datetime.fromtimestamp(
                            file_path.stat().st_ctime
                        ).isoformat(),
                        "format": (
                            file_path.suffix[1:] if file_path.suffix else "unknown"
                        ),
                    }
                )

        return sorted(documents, key=lambda x: x["created"], reverse=True)


class QuickDocumentGenerator:
    """Generador rápido de documentos para casos simples"""

    @staticmethod
    def generate_quick_summary(
        topic: str, content: str, format_type: str = "md"
    ) -> str:
        """Genera resumen rápido"""

        output_dir = Path("exports/quick_docs")
        output_dir.mkdir(parents=True, exist_ok=True)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        if format_type == "md":
            filename = f"summary_{topic.replace(' ', '_')}_{timestamp}.md"
            filepath = output_dir / filename

            md_content = f"""# Resumen: {topic}

**Generado:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

## Contenido

{content}

---
*Generado por Web Scraper PRO*
"""

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(md_content)

        elif format_type == "txt":
            filename = f"summary_{topic.replace(' ', '_')}_{timestamp}.txt"
            filepath = output_dir / filename

            txt_content = f"""RESUMEN: {topic.upper()}
{'='*60}

Generado: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

CONTENIDO:
{content}

Generado por Web Scraper PRO
"""

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(txt_content)

        logger.info(f"📄 Resumen rápido generado: {filepath}")
        return str(filepath)


# Funciones de utilidad
def create_document_generator(config: Dict[str, Any] = None) -> DocumentGenerator:
    """Crea instancia del generador de documentos"""
    return DocumentGenerator(config)


async def quick_document_from_search(
    search_data: Dict[str, Any], format_types: List[str] = None
) -> Dict[str, Any]:
    """Genera documento rápido a partir de datos de búsqueda"""

    generator = create_document_generator()

    return await generator.generate_comprehensive_document(
        search_data,
        document_type="research_report",
        format_types=format_types or ["md", "docx"],
    )


# Función principal para demostración
async def main():
    """Función principal para demostración"""
    print("📄 Iniciando Generador de Documentos...")

    # Datos de ejemplo
    sample_search_data = {
        "topic": "Inteligencia Artificial",
        "synthesis": {
            "summary": "La Inteligencia Artificial es una tecnología transformadora...",
            "key_points": [
                "IA está revolucionando múltiples industrias",
                "Machine Learning es una rama importante de la IA",
                "Considera aspectos éticos importantes",
            ],
            "sources_count": 5,
            "content_quality": 0.85,
            "recommended_actions": [
                "create_comprehensive_document",
                "generate_presentation",
            ],
        },
        "organized_content": {
            "definition": [
                {
                    "title": "Definición de IA",
                    "content": "La Inteligencia Artificial es...",
                    "source": "wikipedia",
                    "url": "https://example.com",
                    "confidence": 0.9,
                }
            ]
        },
    }

    generator = create_document_generator()

    # Generar documento
    result = await generator.generate_comprehensive_document(
        sample_search_data, document_type="research_report", format_types=["md", "docx"]
    )

    if result["status"] == "completed":
        print("✅ Documentos generados exitosamente!")
        for file_info in result["generated_files"]:
            if file_info["status"] == "success":
                print(
                    f"📄 {file_info['format'].upper()}: {file_info.get('path', file_info.get('url', 'N/A'))}"
                )
    else:
        print(f"❌ Error: {result.get('error', 'Error desconocido')}")


if __name__ == "__main__":
    asyncio.run(main())
